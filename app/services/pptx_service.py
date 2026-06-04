from pptx import Presentation

def extract_text_from_pptx(pptx_path):

    presentation = Presentation(pptx_path)

    text = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)